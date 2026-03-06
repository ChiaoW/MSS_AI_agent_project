from ast import pattern
import os
import re
import tempfile
from pathlib import Path
import zipfile
import pandas as pd
from docx import Document
from pptx import Presentation
import extract_msg
import pytesseract
from PIL import Image
import io
import base64
import requests
import logging
from datetime import datetime
from pdf2image import convert_from_path

logger = logging.getLogger(__name__)

class UniversalFileProcessor:
    def __init__(self, temp_dir_base=None, 
                 ocr_api_base="http://localhost:8001/v1",  # 指向 olmOCR Docker (Port 8001)
                 ocr_model="allenai/olmOCR-2-7B-1025"):
        self.temp_dir_base = temp_dir_base or tempfile.gettempdir()
        self.ocr_api_base = ocr_api_base
        self.ocr_model = ocr_model

    def process_directory(self, folder_path: str) -> str:
        folder = Path(folder_path)
        if not folder.exists():
            return ""

        consolidated_text = []
        for file_path in folder.rglob("*"):
            if file_path.is_file():
                if file_path.name.startswith("~") or file_path.name.startswith("."):
                    continue
                
                logger.info(f"Processing: {file_path.name}")
                file_content = self._dispatch_file_handler(file_path)

                if file_content:
                    clean_content = self._advanced_cleaning(file_content)
                    if clean_content.strip():
                        consolidated_text.append(clean_content)

        return "\n".join(consolidated_text)

    def _dispatch_file_handler(self, file_path: Path) -> str:
        """根據副檔名分派處理函式"""
        ext = file_path.suffix.lower()
        
        try:
            if ext == '.txt':
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    content =self._clean_email_body(f.read())
                    # content =f.read()
            elif ext in [".csv", ".xls", ".xlsx"]:
                content =self._read_spreadsheet(file_path)
            elif ext == '.docx':
                content =self._read_word(file_path)
            elif ext == '.pptx':
                content =self._read_pptx(file_path)
            elif ext == '.pdf':
                content =self._read_pdf_with_olmocr(file_path)
            elif ext in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff']:
                content =self._process_image_with_olmocr(file_path)
            elif ext == '.msg':
                content =self._read_msg(file_path)
            elif ext == '.zip':
                content =self._read_zip(file_path)
            else:
                return ""
        except Exception as e:
            return ""
        
        return content
    
    def _read_pdf_with_olmocr(self, path: Path) -> str:
        """
        全 OCR 模式 (All-olmOCR):
        不依賴 pdfplumber，直接將 PDF 的每一頁轉成圖片，交給 olmOCR 處理。
        """
        text_content = []
        try:
            logger.debug(f"Converting PDF to images (this might take a few seconds)")
            # 將 PDF 轉為圖片。dpi=200 是一個很好的平衡點 (解析度夠看清楚小字，且檔案不會太大導致 OOM)
            images = convert_from_path(str(path), dpi=200)
            total_pages = len(images)
            
            for i, img in enumerate(images):
                logger.info(f"[olmOCR] Processing page {i+1}/{total_pages}")
                vlm_text = self._call_olmocr_api(img)
                text_content.append(f"\n--- [Page {i+1}] ---\n{vlm_text}\n")
                
        except Exception as e:
            logger.error(f"PDF to Image Error]: {e}")
            return ""

        return "\n".join(text_content)

    def _process_image_with_olmocr(self, path: Path) -> str:
        try:
            image = Image.open(path)
            return self._call_olmocr_api(image)
        except Exception as e:
            return f"Error processing image: {e}"

    def _call_olmocr_api(self, image: Image) -> str:
        """呼叫本地 Docker 容器的 olmOCR API"""
        buffered = io.BytesIO()
        # 轉換為 RGB 模式，以防 PNG 透明背景導致報錯
        if image.mode != "RGB":
            image = image.convert("RGB")
            
        image.save(buffered, format="JPEG", quality=90)
        img_str = base64.b64encode(buffered.getvalue()).decode("utf-8")
        
        headers = {"Content-Type": "application/json", "Authorization": "Bearer EMPTY"}
        payload = {
            "model": self.ocr_model,
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text", 
                            "text": "Convert this image to markdown. You must strictly use standard markdown table syntax (with '|' and '-') for all tables. Do NOT output any HTML tags such as <table>, <tr>, <th>, or <td> under any circumstances."
                        },
                        {
                            "type": "image_url",
                            "image_url": {"url": f"data:image/jpeg;base64,{img_str}"}
                        }
                    ]
                }
            ],
            "max_tokens": 4096,
            "temperature": 0.1
        }

        try:
            response = requests.post(f"{self.ocr_api_base}/chat/completions", headers=headers, json=payload, timeout=120)
            
            if response.status_code != 200:
                return f"[OCR Error]: Docker responded with {response.status_code} - {response.text}"
                
            result = response.json()
            return result['choices'][0]['message']['content']
        except Exception as e:
            return f"[OCR API Connection Failed]: {e}."
    
    def _advanced_cleaning(self, text: str) -> str:
        """後處理清洗"""
        if not text: return ""
        
        # 清除 Email Header
        pattern = r"(?m)^(?:Cc|Bcc):.*(?:\n(?!\w+:).*)*"
        text = re.sub(pattern, "", text)
        
        # 清除免責聲明
        jp_disclaimer = r"本資料の取扱上の注意.*?東京エレクトロン"
        text = re.sub(jp_disclaimer, "", text, flags=re.DOTALL)
        text = re.sub(r"Internal Use Only", "", text, flags=re.IGNORECASE)
        mss_pattern = r"(?:(?:\*{5,}.*?Notice)|(?:The content of this E-mail may contain)).*?MSSCORPS.*?cooperation\.\s*\*{5,}"
        text = re.sub(mss_pattern, "", text, flags=re.DOTALL | re.IGNORECASE)
        
        def html_table_to_markdown(match):
            logger.debug('Enter table to markdown.')
            html_content = match.group(0)
            try:
                # 使用 pandas 讀取 HTML 表格字串
                dfs = pd.read_html(io.StringIO(html_content))
                if dfs:
                    df = dfs[0].astype(str)
                    # 統一呼叫先前修改過的 Markdown 轉換方法 (需確保已引入或直接在此呼叫 to_markdown)
                    return df.to_markdown(index=False, tablefmt="pipe")
            except Exception as e:
                # 若解析失敗，則保留原樣以避免遺失資料
                logger.warning(f'Fail to turn to markdown, {e}')
                return html_content
            return html_content

        # 使用正規表示式捕捉 <table> 到 </table> 之間的所有內容 (忽略大小寫、跨行比對)
        text = re.sub(r"(?i)<table.*?>.*?</table>", html_table_to_markdown, text, flags=re.DOTALL)
        
        text = re.sub(r'\n{3,}', '\n\n', text)
        return text.strip()
    
    def _df_to_markdown(self, df: pd.DataFrame) -> str:
        if df.empty: return ""
        try:
            df = df.astype(str)
            content_str = " ".join(df.values.flatten())
            if re.search(r"R\s*\d+\s+G\s*\d+\s+B\s*\d+", content_str):
                 return ""
            
            return df.to_markdown(index=False, tablefmt="pipe")
        except:
            return ""

    def _read_spreadsheet(self, path: Path) -> str:
        if path.suffix == ".csv":
            df = pd.read_csv(path)
        else:
            df = pd.read_excel(path)
        return self._df_to_markdown(df)

    def _read_word(self, path: Path) -> str:
        doc = Document(path)
        text = []
        for p in doc.paragraphs: text.append(p.text)
        for table in doc.tables:
            data = [[cell.text.strip() for cell in row.cells] for row in table.rows]
            if data:
                df = pd.DataFrame(data[1:], columns=data[0])
                text.append(self._df_to_markdown(df))
        return "\n".join(text)

    def _read_pptx(self, path: Path) -> str:
        prs = Presentation(path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame: text.append(shape.text_frame.text)
        return "\n".join(text)

    
    def _read_image(self, path: Path) -> str:
        # 使用 OCR 讀取圖片文字
        try:
            image = Image.open(path)
            text = pytesseract.image_to_string(image)
            return f"[OCR Extraction from Image]:\n{text}"
        except Exception as e:
            return f"Error processing image: {e}"

    def _read_msg(self, path: Path) -> str:
        """
        處理 .msg 檔案：
        1. 讀取內文 (Body) 並移除特定的 Disclaimer
        2. 提取附件至暫存資料夾
        3. 遞迴呼叫 process_directory 讀取附件內容
        """
        msg = extract_msg.Message(path)
        
        # 取得內文並執行清理
        raw_body = msg.body
        clean_body = self._clean_email_body(raw_body) if raw_body else ""

        msg_text = [f"Subject: {msg.subject}", f"Body:\n{clean_body}"]
        
        # 處理附件
        if msg.attachments:
            with tempfile.TemporaryDirectory(dir=self.temp_dir_base) as temp_dir:
                for attachment in msg.attachments:
                    # 儲存附件
                    attachment.save(customPath=temp_dir)
                
                # 遞迴讀取暫存資料夾內的檔案
                attachments_content = self.process_directory(temp_dir)
                msg_text.append(attachments_content)
        
        msg.close()
        return "\n".join(msg_text)

    def _clean_email_body(self, text: str) -> str:
        pattern = r"(?:(?:\*{5,}.*?Notice)|(?:The content of this E-mail may contain)).*?MSSCORPS.*?cooperation\.\s*\*{5,}"

        return re.sub(pattern, "", text, flags=re.DOTALL | re.IGNORECASE).strip()

    def _read_zip(self, path: Path) -> str:
        """
        處理 .zip 檔案：
        1. 解壓縮至暫存資料夾
        2. 遞迴呼叫 process_directory
        """
        content = []
        with tempfile.TemporaryDirectory(dir=self.temp_dir_base) as temp_dir:
            try:
                with zipfile.ZipFile(path, 'r') as zip_ref:
                    zip_ref.extractall(temp_dir)
                
                content.append(f"[Contents of ZIP archive {path.name}:]")
                # 遞迴讀取
                content.append(self.process_directory(temp_dir))
            except Exception as e:
                return f"Error processing ZIP: {e}"
        return "\n".join(content)

# 測試用區塊
if __name__ == "__main__":
    processor = UniversalFileProcessor()
    
    # 指定要測試的 Lot IDs
    TARGET_LOT_IDS = [
        "T25122302"
    ]
    
    # 設定基礎輸入路徑與輸出資料夾
    base_input_dir = "./Kang_Yi_Lin_Merged"
    output_dir = "TW_cases_examples"
    
    # 確保輸出資料夾存在，若不存在則建立
    os.makedirs(output_dir, exist_ok=True)
    
    logger.info(f"Starting batch processing for {len(TARGET_LOT_IDS)} lots...\n")

    for lot_id in TARGET_LOT_IDS:
        # 組合該 Lot 的完整路徑 (例如: ./all_cases/T25090401)
        lot_path = os.path.join(base_input_dir, lot_id)
        
        logger.info(f"--- Processing Lot: {lot_id} ---")
        
        if os.path.exists(lot_path):
            # 執行目錄處理
            result_text = processor.process_directory(lot_path)
            
            # 印出前 500 字供快速檢查 (避免字數過多洗版)
            logger.info(f"Preview (First 500 chars):\n{result_text[:500]}\n...")
            
            # 儲存成獨立的 debug 文字檔
            output_filename = os.path.join(output_dir, f"{lot_id}.txt")
            try:
                with open(output_filename, "w", encoding="utf-8") as f:
                    f.write(result_text)
                logger.info(f"-> Successfully saved to: {output_filename}")
            except Exception as e:
                logger.error(f"-> Error saving file {output_filename}: {e}")
                
        else:
            logger.warning(f"-> Path not found: {lot_path}")
    
    logger.info("Batch processing completed.")